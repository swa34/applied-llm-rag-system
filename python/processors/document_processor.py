#!/usr/bin/env python3
"""
Multi-Format Document Processor

Extracts text content from various document formats and converts to markdown.
Supports: PDF, DOCX, PPTX, XLSX, and plain text files.

Production features:
- Robust error handling per file
- Table extraction with markdown formatting
- Slide-by-slide PowerPoint processing
- Speaker notes extraction
- Heading detection from document styles
- Batch processing with summary generation

Author: Scott Allen
"""

import argparse
import json
import os
import re
import sys
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional

import PyPDF2
import requests
from docx import Document
from openpyxl import load_workbook
# Document processing libraries
from pptx import Presentation


class DocumentProcessor:
    """
    Process various document types and extract text content.

    Supports:
    - PowerPoint (PPTX): Slides, tables, speaker notes
    - Word (DOCX): Paragraphs, tables, headings
    - PDF: Page-by-page text extraction
    - Excel (XLSX): Sheet-by-sheet with table formatting
    - Plain text (TXT, MD): Direct reading

    All output is formatted as markdown for downstream RAG processing.
    """

    def __init__(self, output_dir: str = "./output/documents"):
        """
        Initialize the document processor.

        Args:
            output_dir: Directory to save processed markdown files
        """
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.downloads_dir = self.output_dir / "downloads"
        self.downloads_dir.mkdir(exist_ok=True)

    def process_pptx(self, file_path: Path) -> str:
        """
        Extract text from PowerPoint presentation.

        Features:
        - Slide-by-slide extraction with numbering
        - Title detection
        - Bullet point formatting
        - Table extraction with markdown formatting
        - Speaker notes extraction

        Args:
            file_path: Path to the PPTX file

        Returns:
            Extracted text content formatted as markdown
        """
        try:
            prs = Presentation(str(file_path))
            content = []

            # Extract title from metadata
            if prs.core_properties.title:
                content.append(f"# {prs.core_properties.title}\n")

            # Process each slide
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = []

                # Extract slide title
                if slide.shapes.title:
                    title_text = slide.shapes.title.text.strip()
                    if title_text:
                        slide_content.append(f"## Slide {slide_num}: {title_text}\n")
                else:
                    slide_content.append(f"## Slide {slide_num}\n")

                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if text and text not in slide_content:
                            # Format multi-line text as bullet points
                            if '\n' in text:
                                lines = text.split('\n')
                                formatted_lines = []
                                for line in lines:
                                    line = line.strip()
                                    if line:
                                        # Normalize bullet characters
                                        if line.startswith(('•', '-', '*', '▪', '►', '○', '●')):
                                            formatted_lines.append(f"- {line[1:].strip()}")
                                        else:
                                            formatted_lines.append(f"- {line}")
                                text = '\n'.join(formatted_lines)
                            slide_content.append(text)

                # Extract tables
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        table_content = ["\n### Table\n"]

                        # Create markdown table
                        if table.rows:
                            headers = [cell.text.strip() for cell in table.rows[0].cells]
                            if headers:
                                table_content.append("| " + " | ".join(headers) + " |")
                                table_content.append("| " + " | ".join(["---"] * len(headers)) + " |")

                                for row in table.rows[1:]:
                                    row_data = [cell.text.strip() for cell in row.cells]
                                    table_content.append("| " + " | ".join(row_data) + " |")

                        slide_content.append('\n'.join(table_content))

                # Extract speaker notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_content.append(f"\n**Speaker Notes:**\n{notes_text}")

                if slide_content:
                    content.append('\n'.join(slide_content))
                    content.append("\n---\n")

            return '\n'.join(content)

        except Exception as e:
            print(f"  [ERROR] Processing PPTX: {e}")
            return f"[Error processing PowerPoint file: {e}]"

    def process_docx(self, file_path: Path) -> str:
        """
        Extract text from Word document.

        Features:
        - Paragraph extraction
        - Heading detection from styles
        - Table extraction with markdown formatting

        Args:
            file_path: Path to the DOCX file

        Returns:
            Extracted text content formatted as markdown
        """
        try:
            doc = Document(str(file_path))
            content = []

            # Extract paragraphs with heading detection
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    # Detect headings based on style
                    if para.style.name.startswith('Heading'):
                        # Extract heading level (Heading 1, Heading 2, etc.)
                        level_match = re.search(r'\d+', para.style.name)
                        level = int(level_match.group()) if level_match else 1
                        prefix = '#' * level
                        content.append(f"{prefix} {text}\n")
                    else:
                        content.append(text)

            # Extract tables
            for table in doc.tables:
                table_content = ["\n### Table\n"]

                if table.rows:
                    # First row as headers
                    headers = [cell.text.strip() for cell in table.rows[0].cells]

                    if headers:
                        table_content.append("| " + " | ".join(headers) + " |")
                        table_content.append("| " + " | ".join(["---"] * len(headers)) + " |")

                        for row in table.rows[1:]:
                            row_data = [cell.text.strip() for cell in row.cells]
                            table_content.append("| " + " | ".join(row_data) + " |")

                content.append('\n'.join(table_content))

            return '\n\n'.join(content)

        except Exception as e:
            print(f"  [ERROR] Processing DOCX: {e}")
            return f"[Error processing Word document: {e}]"

    def process_pdf(self, file_path: Path) -> str:
        """
        Extract text from PDF document.

        Page-by-page extraction with page number markers.

        Args:
            file_path: Path to the PDF file

        Returns:
            Extracted text content
        """
        try:
            content = []

            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                num_pages = len(pdf_reader.pages)

                for page_num, page in enumerate(pdf_reader.pages, 1):
                    page_text = page.extract_text()
                    if page_text:
                        page_text = page_text.strip()
                        if page_text:
                            content.append(f"### Page {page_num}\n\n{page_text}")

            return '\n\n---\n\n'.join(content)

        except Exception as e:
            print(f"  [ERROR] Processing PDF: {e}")
            return f"[Error processing PDF: {e}]"

    def process_xlsx(self, file_path: Path) -> str:
        """
        Extract text from Excel spreadsheet.

        Sheet-by-sheet extraction with markdown table formatting.

        Args:
            file_path: Path to the XLSX file

        Returns:
            Extracted text content formatted as markdown
        """
        try:
            wb = load_workbook(str(file_path), data_only=True)
            content = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_content = [f"## Sheet: {sheet_name}\n"]

                # Collect rows
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    # Filter out completely empty rows
                    if any(cell is not None for cell in row):
                        rows.append([
                            str(cell) if cell is not None else ""
                            for cell in row
                        ])

                if rows:
                    # First row as headers
                    headers = rows[0]
                    sheet_content.append("| " + " | ".join(headers) + " |")
                    sheet_content.append("| " + " | ".join(["---"] * len(headers)) + " |")

                    # Data rows
                    for row in rows[1:]:
                        # Ensure row has same number of columns as headers
                        while len(row) < len(headers):
                            row.append("")
                        sheet_content.append("| " + " | ".join(row[:len(headers)]) + " |")

                content.append('\n'.join(sheet_content))

            return '\n\n---\n\n'.join(content)

        except Exception as e:
            print(f"  [ERROR] Processing XLSX: {e}")
            return f"[Error processing Excel file: {e}]"

    def download_file(self, url: str, output_path: Path) -> bool:
        """
        Download file from URL.

        Handles Dropbox share links by converting to direct download URLs.

        Args:
            url: URL to download from
            output_path: Path to save the file

        Returns:
            True if successful, False otherwise
        """
        try:
            # Convert Dropbox share links to direct download
            if 'dropbox.com' in url:
                url = url.replace('www.dropbox.com', 'dl.dropboxusercontent.com')
                url = url.replace('dl=0', 'raw=1')

            print(f"  Downloading: {url[:60]}...")

            response = requests.get(url, stream=True, timeout=60)
            response.raise_for_status()

            with open(output_path, 'wb') as f:
                for chunk in response.iter_content(chunk_size=8192):
                    f.write(chunk)

            print(f"  [OK] Downloaded to: {output_path.name}")
            return True

        except Exception as e:
            print(f"  [ERROR] Download failed: {e}")
            return False

    def process_file(self, file_path: Path, source_url: Optional[str] = None) -> Dict:
        """
        Process any supported document type.

        Args:
            file_path: Path to the document file
            source_url: Optional source URL for metadata

        Returns:
            Dictionary with processing results
        """
        file_path = Path(file_path)
        extension = file_path.suffix.lower()

        print(f"Processing {extension.upper()}: {file_path.name}")

        # Extract content based on file type
        content = ""
        if extension in ['.ppt', '.pptx']:
            content = self.process_pptx(file_path)
        elif extension in ['.doc', '.docx']:
            content = self.process_docx(file_path)
        elif extension == '.pdf':
            content = self.process_pdf(file_path)
        elif extension in ['.xls', '.xlsx']:
            content = self.process_xlsx(file_path)
        elif extension in ['.txt', '.md']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
        else:
            content = f"[Unsupported file type: {extension}]"

        # Create markdown document with metadata
        title = file_path.stem.replace('-', ' ').replace('_', ' ').title()

        markdown = f"""# {title}

**Source:** {source_url or 'Local file'}
**Type:** {extension[1:].upper()}
**Processed:** {datetime.now().isoformat()}

---

{content}
"""

        # Save markdown file
        output_file = self.output_dir / f"{file_path.stem}.md"

        # Handle duplicates
        counter = 1
        while output_file.exists():
            output_file = self.output_dir / f"{file_path.stem}_{counter}.md"
            counter += 1

        output_file.write_text(markdown, encoding='utf-8')

        print(f"  [OK] Saved: {output_file.name}")

        return {
            'title': title,
            'file_path': str(file_path),
            'output_path': str(output_file),
            'extension': extension,
            'source_url': source_url,
            'content_length': len(content),
            'status': 'success'
        }

    def process_url(self, url: str) -> Dict:
        """
        Download and process a file from URL.

        Args:
            url: URL to download and process

        Returns:
            Dictionary with processing results
        """
        # Extract filename from URL
        url_parts = url.split('/')
        filename_part = url_parts[-1] if url_parts else 'document'
        filename = filename_part.split('?')[0]

        # Ensure valid extension
        if not any(filename.lower().endswith(ext) for ext in ['.pdf', '.docx', '.pptx', '.xlsx', '.txt']):
            filename += '.pdf'  # Default assumption

        # Download file
        download_path = self.downloads_dir / filename

        if self.download_file(url, download_path):
            return self.process_file(download_path, source_url=url)
        else:
            return {
                'url': url,
                'status': 'failed',
                'error': 'Download failed'
            }

    def process_batch(self, items: List[str]) -> Dict:
        """
        Process multiple files or URLs.

        Args:
            items: List of file paths or URLs

        Returns:
            Dictionary with processing summary
        """
        results = []
        errors = []

        print(f"\n{'=' * 60}")
        print(f"Processing {len(items)} items")
        print(f"{'=' * 60}\n")

        for i, item in enumerate(items, 1):
            print(f"\n[{i}/{len(items)}]")

            try:
                if item.startswith('http'):
                    result = self.process_url(item)
                else:
                    result = self.process_file(Path(item))

                if result.get('status') == 'success':
                    results.append(result)
                else:
                    errors.append(result)

            except Exception as e:
                errors.append({
                    'item': item,
                    'error': str(e),
                    'status': 'failed'
                })
                print(f"  [ERROR] {e}")

        # Save processing summary
        summary = {
            'processed_at': datetime.now().isoformat(),
            'total_items': len(items),
            'successful': len(results),
            'failed': len(errors),
            'results': results,
            'errors': errors
        }

        summary_path = self.output_dir / 'processing_summary.json'
        with open(summary_path, 'w', encoding='utf-8') as f:
            json.dump(summary, f, indent=2)

        print(f"\n{'=' * 60}")
        print(f"Processing Complete")
        print(f"{'=' * 60}")
        print(f"  Successful: {len(results)}/{len(items)}")
        print(f"  Failed: {len(errors)}")
        print(f"  Output: {self.output_dir}")
        print(f"  Summary: {summary_path}")

        return summary


def main():
    """Command-line interface."""
    parser = argparse.ArgumentParser(
        description='Process documents and convert to markdown for RAG ingestion'
    )
    parser.add_argument(
        'items',
        nargs='+',
        help='File paths or URLs to process'
    )
    parser.add_argument(
        '--output-dir',
        default='./output/documents',
        help='Output directory for markdown files'
    )
    parser.add_argument(
        '--test',
        action='store_true',
        help='Test mode - process first item only'
    )

    args = parser.parse_args()

    # Create processor
    processor = DocumentProcessor(output_dir=args.output_dir)

    # Process items
    items = args.items[:1] if args.test else args.items
    if args.test:
        print("Test mode: Processing first item only")

    summary = processor.process_batch(items)

    # Exit with error code if any failed
    sys.exit(1 if summary['failed'] > 0 else 0)


if __name__ == '__main__':
    main()
